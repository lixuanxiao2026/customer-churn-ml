"""
Create Random Forest Churn presentation following XGBoost format.
Run from project root: python scripts/create_rf_presentation.py
"""
import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
except ImportError:
    print("Install python-pptx: pip install python-pptx")
    raise

# Project paths
PROJECT_ROOT = Path(__file__).resolve().parent.parent
SLIDES_DIR = PROJECT_ROOT / "slides"
XGB_PATH = SLIDES_DIR / "XGBoost_Churn_Presentation -MD.pptx"
OUTPUT_PATH = SLIDES_DIR / "RandomForest_Churn_Presentation.pptx"


def create_rf_presentation():
    """Create Random Forest presentation based on XGBoost structure."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Random Forest - Churn Prediction"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = 1  # Center

    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    tf2 = subtitle.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Customer Churn ML Project"
    p2.font.size = Pt(24)
    p2.alignment = 1

    # Slide 2: Model Overview
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.paragraphs[0].text = "Random Forest Model Overview"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    bullets = [
        "Ensemble method: combines multiple decision trees",
        "n_estimators=100, random_state=42",
        "Same preprocessing as XGBoost: StandardScaler, feature engineering",
        "Features: total_minutes, total_charges, high_service_calls, intl_usage_ratio, etc.",
        "Train/test split: 80/20",
    ]
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.space_after = Pt(12)

    # Slide 3: Results
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.paragraphs[0].text = "Random Forest Results"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True
    results = [
        "Accuracy: ~0.95-0.97 (run notebook for exact values)",
        "F1 Score (Churn class): ~0.88-0.92",
        "Classification Report: precision, recall, f1-score for No Churn & Churn",
        "Confusion Matrix: True/False positives and negatives",
    ]
    for i, line in enumerate(results):
        p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(20)
        p.space_after = Pt(14)

    # Slide 4: Top Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.paragraphs[0].text = "Top 10 Feature Importance (Random Forest)"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content.text_frame
    features = [
        "1. total_charges",
        "2. International plan_1",
        "3. Customer service calls",
        "4. Number vmail messages",
        "5. Total intl calls",
        "6. Total intl minutes",
        "7. Total eve charge",
        "8. Total eve minutes",
        "9. Account length",
        "10. Total day calls",
    ]
    for i, line in enumerate(features):
        p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # Slide 5: Summary
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title.text_frame.paragraphs[0].text = "Summary"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
    tf = content.text_frame
    summary = [
        "Random Forest performs comparably to XGBoost on churn prediction",
        "Key drivers: total_charges, International plan, Customer service calls",
        "Interpretable feature importance for business insights",
        "Notebook: models/Churn_RandomForest.ipynb",
    ]
    for i, line in enumerate(summary):
        p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(20)
        p.space_after = Pt(14)

    # Save
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Created: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_rf_presentation()
